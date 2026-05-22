import os
import csv
import traceback
import importlib
from logger import Logger

SUPPORTED_EXTENSIONS = {
    ".doc", ".docx",
    ".xls", ".xlsx", ".csv",
    ".ppt", ".pptx",
    ".txt", ".md",
    ".pdf",
    ".wps", ".et", ".dps"
}

# Windows 系统目录黑名单（路径前缀，不区分大小写）
# 这些目录属于系统/程序安装区域，不含用户数据，跳过可大幅减少无效扫描
WINDOWS_SYSTEM_DIRS = {
    "windows",
    "program files",
    "program files (x86)",
    "programdata",
    "$recycle.bin",
    "system volume information",
}


def _is_windows_system_dir(parent_path, dir_name):
    """判断目录是否为 Windows 系统目录（仅在驱动器根下才过滤）

    只过滤驱动器根目录下的系统目录（如 C:\\Windows、C:\\Program Files），
    避免误过滤用户自建的同名子文件夹。
    """
    # 判断 parent_path 是否为驱动器根（如 C:\ 或 C:）
    normalized_parent = os.path.normcase(os.path.normpath(parent_path))
    # 驱动器根格式：X:\ 或 X:
    drive, tail = os.path.splitdrive(normalized_parent)
    if not drive:
        return False
    tail_stripped = tail.strip(os.sep)
    if tail_stripped:  # 非驱动器根（还有子路径）
        return False
    return dir_name.lower() in WINDOWS_SYSTEM_DIRS

# ============================================================
# 模块级 import 缓存：避免每次调用方法时重复 import 查找
# ============================================================
_cached_imports = {}

# 文本提取上限：防止单行超长或极端文件占用过多内存
MAX_TEXT_CHARS = 10_000_000  # 约 10MB 文本

# OLE 流读取上限：避免一次性加载超大二进制流
MAX_OLE_STREAM_BYTES = 2 * 1024 * 1024  # 2MB


def _lazy_import(module_name):
    """延迟导入并缓存，仅在首次使用时真正 import

    使用 importlib.import_module 替代 __import__，
    对 fitz (PyMuPDF) 等包名与模块名不一致的库兼容性更好。
    """
    if module_name not in _cached_imports:
        try:
            _cached_imports[module_name] = importlib.import_module(module_name)
        except ImportError:
            _cached_imports[module_name] = None
    return _cached_imports[module_name]


class FileScanner:
    def __init__(self, max_file_size_bytes):
        self.max_file_size_bytes = max_file_size_bytes
        self.logger = Logger()

    def scan_directory(self, root_dir):
        if not os.path.exists(root_dir):
            self.logger.error(f"扫描目录不存在: {root_dir}")
            return []

        target_files = []
        for dirpath, dirnames, filenames in os.walk(root_dir):
            # 跳过 Windows 系统目录：原地修改 dirnames 阻止 os.walk 递归进入
            dirnames[:] = [
                d for d in dirnames
                if not _is_windows_system_dir(dirpath, d)
            ]
            for filename in filenames:
                ext = os.path.splitext(filename)[1].lower()
                if ext in SUPPORTED_EXTENSIONS:
                    full_path = os.path.join(dirpath, filename)
                    target_files.append(full_path)

        self.logger.info(f"扫描到 {len(target_files)} 个待检测文件")
        return target_files

    def get_file_size_mb(self, file_path):
        try:
            size_bytes = os.path.getsize(file_path)
            return size_bytes / (1024 * 1024)
        except OSError:
            return -1

    def is_oversized(self, file_path):
        size_mb = self.get_file_size_mb(file_path)
        if size_mb < 0:
            self.logger.warning(f"无法获取文件大小: {file_path}")
            return True, -1
        max_size_mb = self.max_file_size_bytes / (1024 * 1024)
        is_over = size_mb > max_size_mb
        if is_over:
            self.logger.info(f"文件超过大小阈值 ({size_mb:.2f}MB > {max_size_mb:.0f}MB)，跳过: {file_path}")
        return is_over, size_mb

    def extract_text(self, file_path):
        """提取文件全部文本（原有接口，保持兼容）"""
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == ".txt":
                return self._read_plain_text(file_path)
            elif ext == ".docx":
                return self._read_docx(file_path)
            elif ext == ".doc":
                return self._read_doc(file_path)
            elif ext == ".xlsx":
                return self._read_xlsx(file_path)
            elif ext == ".xls":
                return self._read_xls(file_path)
            elif ext == ".csv":
                return self._read_csv(file_path)
            elif ext == ".pptx":
                return self._read_pptx(file_path)
            elif ext == ".ppt":
                return self._read_ppt(file_path)
            elif ext == ".pdf":
                return self._read_pdf(file_path)
            elif ext == ".wps":
                return self._read_wps(file_path)
            elif ext == ".et":
                return self._read_xls(file_path)
            elif ext == ".dps":
                return self._read_ppt(file_path)
            else:
                return ""
        except Exception as e:
            self.logger.warning(f"文件解析失败 ({ext}): {file_path} - {e}")
            self.logger.debug(traceback.format_exc())
            return ""

    def extract_and_match(self, file_path, matcher):
        """边提取边匹配：增量提取文本并实时检测，达到阈值立即停止提取

        相比 extract_text + scan_text 的两阶段模式，此方法可以在提取过程中
        一旦发现达到匹配阈值就立即终止，避免提取剩余无用文本。

        Args:
            file_path: 文件路径
            matcher: Matcher 实例

        Returns:
            (is_sensitive, details, size_mb):
                is_sensitive: 是否敏感
                details: 匹配详情列表
                size_mb: 文件大小(MB)
        """
        ext = os.path.splitext(file_path)[1].lower()

        # 先检查文件大小
        is_over, size_mb = self.is_oversized(file_path)
        if is_over:
            return False, [], size_mb

        try:
            if ext in (".txt", ".md", ".log"):
                return self._read_plain_text_incremental(file_path, matcher, size_mb)
            elif ext == ".docx":
                return self._read_docx_incremental(file_path, matcher, size_mb)
            elif ext == ".xlsx":
                return self._read_xlsx_incremental(file_path, matcher, size_mb)
            elif ext == ".xls":
                return self._read_xls_incremental(file_path, matcher, size_mb)
            elif ext == ".csv":
                return self._read_csv_incremental(file_path, matcher, size_mb)
            elif ext == ".pptx":
                return self._read_pptx_incremental(file_path, matcher, size_mb)
            elif ext == ".pdf":
                return self._read_pdf_incremental(file_path, matcher, size_mb)
            else:
                # 其他格式回退到全量提取+匹配
                text = self.extract_text(file_path)
                is_sensitive, details = matcher.scan_text(text, file_path)
                return is_sensitive, details, size_mb
        except Exception as e:
            self.logger.warning(f"文件解析失败 ({ext}): {file_path} - {e}")
            self.logger.debug(traceback.format_exc())
            return False, [], size_mb

    # ----------------------------------------------------------
    #  纯文本提取方法（保持原有逻辑）
    # ----------------------------------------------------------

    def _read_plain_text(self, file_path):
        content = []
        total_chars = 0
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                for i, line in enumerate(f):
                    if i >= 10000 or total_chars >= MAX_TEXT_CHARS:
                        break
                    content.append(line)
                    total_chars += len(line)
        except (IOError, PermissionError) as e:
            self.logger.warning(f"文本文件读取失败: {file_path} - {e}")
            return ""
        return "".join(content)

    def _read_docx(self, file_path):
        docx_mod = _lazy_import("docx")
        if docx_mod is None:
            self.logger.error("缺少 python-docx 库，无法解析 .docx 文件")
            return ""
        try:
            from docx import Document
            doc = Document(file_path)
            content = []
            line_count = 0
            total_chars = 0
            for para in doc.paragraphs:
                if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                    break
                content.append(para.text)
                line_count += 1
                total_chars += len(para.text)
            for table in doc.tables:
                if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                    break
                for row in table.rows:
                    if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                        break
                    row_text = " ".join(cell.text for cell in row.cells)
                    content.append(row_text)
                    line_count += 1
                    total_chars += len(row_text)
            return "\n".join(content)
        except Exception as e:
            self.logger.warning(f"docx 解析失败: {file_path} - {e}")
            return ""

    def _read_doc(self, file_path):
        olefile_mod = _lazy_import("olefile")
        if olefile_mod is None:
            self.logger.error("缺少 olefile 库，无法解析 .doc 文件")
            return ""
        try:
            import olefile
            if not olefile.isOleFile(file_path):
                self.logger.warning(f"不是有效的 OLE 文件: {file_path}")
                return ""
            ole = olefile.OleFileIO(file_path)
            try:
                word_stream = ole.openstream("WordDocument")
                # 仅读取需要的字节数，避免大文件一次性加载全部数据
                data = word_stream.read(20002)
                word_stream.close()

                text_parts = []
                text_mode = False
                for i in range(0, len(data) - 1, 2):
                    if i >= 20000:
                        break
                    char_code = data[i] | (data[i + 1] << 8)
                    if char_code == 0x0D:
                        text_parts.append("\n")
                        text_mode = False
                    elif 0x20 <= char_code <= 0xFFFF and char_code != 0xFEFF:
                        if char_code <= 0xFFFF:
                            text_parts.append(chr(char_code))
                            text_mode = True

                return "".join(text_parts)
            finally:
                ole.close()
        except Exception as e:
            self.logger.warning(f"doc 解析失败: {file_path} - {e}")
            return ""

    @staticmethod
    def _cell_to_str(value):
        if isinstance(value, float):
            if value == int(value):
                return str(int(value))
        return str(value)

    def _read_xlsx(self, file_path):
        openpyxl_mod = _lazy_import("openpyxl")
        if openpyxl_mod is None:
            self.logger.error("缺少 openpyxl 库，无法解析 .xlsx 文件")
            return ""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            try:
                content = []
                line_count = 0
                total_chars = 0
                for sheet in wb.worksheets:
                    if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                        break
                    for row in sheet.iter_rows(values_only=True):
                        if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                            break
                        row_text = " ".join(self._cell_to_str(cell) for cell in row if cell is not None)
                        if row_text.strip():
                            content.append(row_text)
                            line_count += 1
                            total_chars += len(row_text)
                return "\n".join(content)
            finally:
                wb.close()
        except Exception as e:
            self.logger.warning(f"xlsx 解析失败: {file_path} - {e}")
            return ""

    def _read_xls(self, file_path):
        xlrd_mod = _lazy_import("xlrd")
        if xlrd_mod is None:
            self.logger.error("缺少 xlrd 库，无法解析 .xls/.et 文件")
            return ""
        try:
            import xlrd
            wb = xlrd.open_workbook(file_path)
            content = []
            line_count = 0
            total_chars = 0
            for sheet_idx in range(wb.nsheets):
                if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                    break
                sheet = wb.sheet_by_index(sheet_idx)
                for row_idx in range(sheet.nrows):
                    if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                        break
                    row_values = sheet.row_values(row_idx)
                    row_text = " ".join(self._cell_to_str(v) for v in row_values if v != "")
                    if row_text.strip():
                        content.append(row_text)
                        line_count += 1
                        total_chars += len(row_text)
            return "\n".join(content)
        except Exception as e:
            self.logger.warning(f"xls 解析失败: {file_path} - {e}")
            return ""

    def _read_csv(self, file_path):
        content = []
        total_chars = 0
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                for i, row in enumerate(reader):
                    if i >= 10000 or total_chars >= MAX_TEXT_CHARS:
                        break
                    row_text = " ".join(row)
                    if row_text.strip():
                        content.append(row_text)
                        total_chars += len(row_text)
        except (IOError, PermissionError) as e:
            self.logger.warning(f"CSV 文件读取失败: {file_path} - {e}")
            return ""
        return "\n".join(content)

    def _read_pptx(self, file_path):
        pptx_mod = _lazy_import("pptx")
        if pptx_mod is None:
            self.logger.error("缺少 python-pptx 库，无法解析 .pptx 文件")
            return ""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            content = []
            line_count = 0
            total_chars = 0
            for slide in prs.slides:
                if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                    break
                for shape in slide.shapes:
                    if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                        break
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                                break
                            text = para.text.strip()
                            if text:
                                content.append(text)
                                line_count += 1
                                total_chars += len(text)
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                                break
                            row_text = " ".join(cell.text for cell in row.cells)
                            if row_text.strip():
                                content.append(row_text)
                                line_count += 1
                                total_chars += len(row_text)
            return "\n".join(content)
        except Exception as e:
            self.logger.warning(f"pptx 解析失败: {file_path} - {e}")
            return ""

    def _read_ppt(self, file_path):
        olefile_mod = _lazy_import("olefile")
        if olefile_mod is None:
            self.logger.error("缺少 olefile 库，无法解析 .ppt/.dps 文件")
            return ""
        try:
            import olefile
            if not olefile.isOleFile(file_path):
                self.logger.warning(f"不是有效的 OLE 文件: {file_path}")
                return ""
            ole = olefile.OleFileIO(file_path)
            try:
                content = []
                line_count = 0
                for stream_name in ole.listdir():
                    if line_count >= 10000:
                        break
                    name = "/".join(stream_name).lower()
                    if "text" in name or "word" in name:
                        try:
                            stream = ole.openstream("/".join(stream_name))
                            data = stream.read(MAX_OLE_STREAM_BYTES)
                            stream.close()
                            text = data.decode("utf-8", errors="replace")
                            for line in text.split("\n"):
                                if line_count >= 10000:
                                    break
                                if line.strip():
                                    content.append(line.strip())
                                    line_count += 1
                        except Exception:
                            continue
                return "\n".join(content)
            finally:
                ole.close()
        except Exception as e:
            self.logger.warning(f"ppt 解析失败: {file_path} - {e}")
            return ""

    def _read_pdf(self, file_path):
        """PDF 解析：优先 PyMuPDF → pdfplumber → pypdf/PyPDF2，逐级回退

        PyMuPDF (fitz): C 扩展，速度最快
        pdfplumber: 纯 Python，解析精度高但慢
        pypdf/PyPDF2: 纯 Python，速度适中，作为最终回退
        """
        # 优先尝试 PyMuPDF
        fitz_mod = _lazy_import("fitz")
        if fitz_mod is not None:
            return self._read_pdf_fitz(file_path)

        # 其次尝试 pdfplumber
        pdfplumber_mod = _lazy_import("pdfplumber")
        if pdfplumber_mod is not None:
            return self._read_pdf_pdfplumber(file_path)

        # 最后尝试 pypdf (新包名) 或 PyPDF2 (旧包名)
        pypdf_mod = _lazy_import("pypdf")
        if pypdf_mod is not None:
            return self._read_pdf_pypdf(file_path)

        pypdf2_mod = _lazy_import("PyPDF2")
        if pypdf2_mod is not None:
            return self._read_pdf_pypdf2(file_path)

        self.logger.error(
            "缺少 PDF 解析库，无法解析 .pdf 文件。"
            "请安装以下任一库: pip install PyMuPDF / pip install pdfplumber / pip install pypdf"
        )
        return ""

    def _read_pdf_fitz(self, file_path):
        """使用 PyMuPDF/fitz 解析 PDF（高性能，C 扩展）"""
        try:
            import fitz
            doc = fitz.open(file_path)
            try:
                content = []
                line_count = 0
                total_chars = 0
                for page in doc:
                    if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                        break
                    text = page.get_text()
                    if text:
                        for line in text.split("\n"):
                            if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                                break
                            if line.strip():
                                content.append(line.strip())
                                line_count += 1
                                total_chars += len(line)
                return "\n".join(content)
            finally:
                doc.close()
        except Exception as e:
            self.logger.warning(f"PyMuPDF 解析失败，回退到 pdfplumber: {file_path} - {e}")
            return self._read_pdf_pdfplumber(file_path)

    def _read_pdf_pdfplumber(self, file_path):
        """使用 pdfplumber 解析 PDF（回退方案，纯 Python）"""
        try:
            import pdfplumber
            content = []
            line_count = 0
            total_chars = 0
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                        break
                    text = page.extract_text()
                    if text:
                        for line in text.split("\n"):
                            if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                                break
                            if line.strip():
                                content.append(line.strip())
                                line_count += 1
                                total_chars += len(line)
            return "\n".join(content)
        except Exception as e:
            self.logger.warning(f"PDF 解析失败: {file_path} - {e}")
            return ""

    def _read_wps(self, file_path):
        olefile_mod = _lazy_import("olefile")
        if olefile_mod is None:
            self.logger.error("缺少 olefile 库，无法解析 .wps 文件")
            return ""
        try:
            import olefile
            if not olefile.isOleFile(file_path):
                self.logger.warning(f"不是有效的 OLE 文件 (wps): {file_path}")
                return ""
            ole = olefile.OleFileIO(file_path)
            try:
                text_parts = []
                line_count = 0
                for stream_name in ole.listdir():
                    if line_count >= 10000:
                        break
                    name = "/".join(stream_name).lower()
                    if "text" in name or "content" in name or "worddocument" in name:
                        try:
                            stream = ole.openstream("/".join(stream_name))
                            data = stream.read(MAX_OLE_STREAM_BYTES)
                            stream.close()
                            text = data.decode("utf-8", errors="replace")
                            for line in text.split("\n"):
                                if line_count >= 10000:
                                    break
                                if line.strip():
                                    text_parts.append(line.strip())
                                    line_count += 1
                        except Exception:
                            continue
                if not text_parts:
                    stream = ole.openstream("WordDocument")
                    data = stream.read(20002)
                    stream.close()
                    for i in range(0, min(len(data), 20000), 2):
                        char_code = data[i] | (data[i + 1] << 8)
                        if 0x20 <= char_code <= 0xFFFF and char_code != 0xFEFF:
                            text_parts.append(chr(char_code))
                            if char_code in (0x0D, 0x0A):
                                line_count += 1
                return "".join(text_parts)
            finally:
                ole.close()
        except Exception as e:
            self.logger.warning(f"wps 解析失败: {file_path} - {e}")
            return ""

    # ----------------------------------------------------------
    #  增量提取 + 匹配方法：边提取边检测，达到阈值立即停止
    # ----------------------------------------------------------

    def _read_plain_text_incremental(self, file_path, matcher, size_mb):
        """纯文本增量提取+匹配"""
        accumulated = None
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                chunk_lines = []
                chunk_size = 0
                CHUNK_THRESHOLD = 500  # 每 500 行做一次匹配检测

                for i, line in enumerate(f):
                    if i >= 10000:
                        break
                    chunk_lines.append(line)
                    chunk_size += 1

                    if chunk_size >= CHUNK_THRESHOLD:
                        chunk_text = "".join(chunk_lines)
                        should_stop, accumulated = matcher.scan_text_incremental(
                            chunk_text, accumulated
                        )
                        chunk_lines = []
                        chunk_size = 0

                        if should_stop:
                            is_sensitive, details = matcher.get_match_details(accumulated)
                            return is_sensitive, details, size_mb

                # 处理剩余行
                if chunk_lines:
                    chunk_text = "".join(chunk_lines)
                    should_stop, accumulated = matcher.scan_text_incremental(
                        chunk_text, accumulated
                    )

        except (IOError, PermissionError) as e:
            self.logger.warning(f"文本文件读取失败: {file_path} - {e}")

        is_sensitive, details = matcher.get_match_details(accumulated or {})
        return is_sensitive, details, size_mb

    def _read_docx_incremental(self, file_path, matcher, size_mb):
        """docx 增量提取+匹配"""
        docx_mod = _lazy_import("docx")
        if docx_mod is None:
            return False, [], size_mb
        try:
            from docx import Document
            doc = Document(file_path)
            accumulated = None
            chunk_parts = []
            chunk_size = 0
            CHUNK_THRESHOLD = 200

            for para in doc.paragraphs:
                chunk_parts.append(para.text)
                chunk_size += 1
                if chunk_size >= CHUNK_THRESHOLD:
                    chunk_text = "\n".join(chunk_parts)
                    should_stop, accumulated = matcher.scan_text_incremental(
                        chunk_text, accumulated
                    )
                    chunk_parts = []
                    chunk_size = 0
                    if should_stop:
                        is_sensitive, details = matcher.get_match_details(accumulated)
                        return is_sensitive, details, size_mb

            for table in doc.tables:
                for row in table.rows:
                    row_text = " ".join(cell.text for cell in row.cells)
                    chunk_parts.append(row_text)
                    chunk_size += 1
                    if chunk_size >= CHUNK_THRESHOLD:
                        chunk_text = "\n".join(chunk_parts)
                        should_stop, accumulated = matcher.scan_text_incremental(
                            chunk_text, accumulated
                        )
                        chunk_parts = []
                        chunk_size = 0
                        if should_stop:
                            is_sensitive, details = matcher.get_match_details(accumulated)
                            return is_sensitive, details, size_mb

            if chunk_parts:
                chunk_text = "\n".join(chunk_parts)
                should_stop, accumulated = matcher.scan_text_incremental(
                    chunk_text, accumulated
                )

            is_sensitive, details = matcher.get_match_details(accumulated or {})
            return is_sensitive, details, size_mb
        except Exception as e:
            self.logger.warning(f"docx 增量解析失败: {file_path} - {e}")
            return False, [], size_mb

    def _read_xlsx_incremental(self, file_path, matcher, size_mb):
        """xlsx 增量提取+匹配"""
        openpyxl_mod = _lazy_import("openpyxl")
        if openpyxl_mod is None:
            return False, [], size_mb
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            try:
                accumulated = None
                chunk_parts = []
                chunk_size = 0
                CHUNK_THRESHOLD = 100
                line_count = 0

                for sheet in wb.worksheets:
                    if line_count >= 10000:
                        break
                    for row in sheet.iter_rows(values_only=True):
                        if line_count >= 10000:
                            break
                        row_text = " ".join(self._cell_to_str(cell) for cell in row if cell is not None)
                        if row_text.strip():
                            chunk_parts.append(row_text)
                            chunk_size += 1
                            line_count += 1
                            if chunk_size >= CHUNK_THRESHOLD:
                                chunk_text = "\n".join(chunk_parts)
                                should_stop, accumulated = matcher.scan_text_incremental(
                                    chunk_text, accumulated
                                )
                                chunk_parts = []
                                chunk_size = 0
                                if should_stop:
                                    is_sensitive, details = matcher.get_match_details(accumulated)
                                    return is_sensitive, details, size_mb

                if chunk_parts:
                    chunk_text = "\n".join(chunk_parts)
                    should_stop, accumulated = matcher.scan_text_incremental(
                        chunk_text, accumulated
                    )

                is_sensitive, details = matcher.get_match_details(accumulated or {})
                return is_sensitive, details, size_mb
            finally:
                wb.close()
        except Exception as e:
            self.logger.warning(f"xlsx 增量解析失败: {file_path} - {e}")
            return False, [], size_mb

    def _read_xls_incremental(self, file_path, matcher, size_mb):
        """xls 增量提取+匹配"""
        xlrd_mod = _lazy_import("xlrd")
        if xlrd_mod is None:
            return False, [], size_mb
        try:
            import xlrd
            wb = xlrd.open_workbook(file_path)
            accumulated = None
            chunk_parts = []
            chunk_size = 0
            CHUNK_THRESHOLD = 100
            line_count = 0

            for sheet_idx in range(wb.nsheets):
                if line_count >= 10000:
                    break
                sheet = wb.sheet_by_index(sheet_idx)
                for row_idx in range(sheet.nrows):
                    if line_count >= 10000:
                        break
                    row_values = sheet.row_values(row_idx)
                    row_text = " ".join(self._cell_to_str(v) for v in row_values if v != "")
                    if row_text.strip():
                        chunk_parts.append(row_text)
                        chunk_size += 1
                        line_count += 1
                        if chunk_size >= CHUNK_THRESHOLD:
                            chunk_text = "\n".join(chunk_parts)
                            should_stop, accumulated = matcher.scan_text_incremental(
                                chunk_text, accumulated
                            )
                            chunk_parts = []
                            chunk_size = 0
                            if should_stop:
                                is_sensitive, details = matcher.get_match_details(accumulated)
                                return is_sensitive, details, size_mb

            if chunk_parts:
                chunk_text = "\n".join(chunk_parts)
                should_stop, accumulated = matcher.scan_text_incremental(
                    chunk_text, accumulated
                )

            is_sensitive, details = matcher.get_match_details(accumulated or {})
            return is_sensitive, details, size_mb
        except Exception as e:
            self.logger.warning(f"xls 增量解析失败: {file_path} - {e}")
            return False, [], size_mb

    def _read_csv_incremental(self, file_path, matcher, size_mb):
        """csv 增量提取+匹配"""
        accumulated = None
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                reader = csv.reader(f)
                chunk_parts = []
                chunk_size = 0
                CHUNK_THRESHOLD = 100

                for i, row in enumerate(reader):
                    if i >= 10000:
                        break
                    row_text = " ".join(row)
                    if row_text.strip():
                        chunk_parts.append(row_text)
                        chunk_size += 1
                        if chunk_size >= CHUNK_THRESHOLD:
                            chunk_text = "\n".join(chunk_parts)
                            should_stop, accumulated = matcher.scan_text_incremental(
                                chunk_text, accumulated
                            )
                            chunk_parts = []
                            chunk_size = 0
                            if should_stop:
                                is_sensitive, details = matcher.get_match_details(accumulated)
                                return is_sensitive, details, size_mb

            if chunk_parts:
                chunk_text = "\n".join(chunk_parts)
                should_stop, accumulated = matcher.scan_text_incremental(
                    chunk_text, accumulated
                )

        except (IOError, PermissionError) as e:
            self.logger.warning(f"CSV 文件读取失败: {file_path} - {e}")

        is_sensitive, details = matcher.get_match_details(accumulated or {})
        return is_sensitive, details, size_mb

    def _read_pptx_incremental(self, file_path, matcher, size_mb):
        """pptx 增量提取+匹配"""
        pptx_mod = _lazy_import("pptx")
        if pptx_mod is None:
            return False, [], size_mb
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            accumulated = None
            chunk_parts = []
            chunk_size = 0
            CHUNK_THRESHOLD = 50
            line_count = 0

            for slide in prs.slides:
                if line_count >= 10000:
                    break
                for shape in slide.shapes:
                    if line_count >= 10000:
                        break
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if line_count >= 10000:
                                break
                            text = para.text.strip()
                            if text:
                                chunk_parts.append(text)
                                chunk_size += 1
                                line_count += 1
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            if line_count >= 10000:
                                break
                            row_text = " ".join(cell.text for cell in row.cells)
                            if row_text.strip():
                                chunk_parts.append(row_text)
                                chunk_size += 1
                                line_count += 1

                    if chunk_size >= CHUNK_THRESHOLD:
                        chunk_text = "\n".join(chunk_parts)
                        should_stop, accumulated = matcher.scan_text_incremental(
                            chunk_text, accumulated
                        )
                        chunk_parts = []
                        chunk_size = 0
                        if should_stop:
                            is_sensitive, details = matcher.get_match_details(accumulated)
                            return is_sensitive, details, size_mb

            if chunk_parts:
                chunk_text = "\n".join(chunk_parts)
                should_stop, accumulated = matcher.scan_text_incremental(
                    chunk_text, accumulated
                )

            is_sensitive, details = matcher.get_match_details(accumulated or {})
            return is_sensitive, details, size_mb
        except Exception as e:
            self.logger.warning(f"pptx 增量解析失败: {file_path} - {e}")
            return False, [], size_mb

    def _read_pdf_incremental(self, file_path, matcher, size_mb):
        """PDF 增量提取+匹配：优先 PyMuPDF → pdfplumber → pypdf/PyPDF2"""
        fitz_mod = _lazy_import("fitz")
        if fitz_mod is not None:
            return self._read_pdf_fitz_incremental(file_path, matcher, size_mb)

        pdfplumber_mod = _lazy_import("pdfplumber")
        if pdfplumber_mod is not None:
            return self._read_pdf_pdfplumber_incremental(file_path, matcher, size_mb)

        # pypdf/PyPDF2 增量提取
        pypdf_mod = _lazy_import("pypdf")
        if pypdf_mod is not None:
            return self._read_pdf_pypdf_incremental(file_path, matcher, size_mb)

        pypdf2_mod = _lazy_import("PyPDF2")
        if pypdf2_mod is not None:
            return self._read_pdf_pypdf2_incremental(file_path, matcher, size_mb)

        self.logger.error(
            "缺少 PDF 解析库，无法解析 .pdf 文件。"
            "请安装以下任一库: pip install PyMuPDF / pip install pdfplumber / pip install pypdf"
        )
        return False, [], size_mb

    def _read_pdf_fitz_incremental(self, file_path, matcher, size_mb):
        """使用 PyMuPDF 增量提取+匹配"""
        try:
            import fitz
            doc = fitz.open(file_path)
            try:
                accumulated = None
                chunk_parts = []
                chunk_size = 0
                CHUNK_THRESHOLD = 10  # PDF 每页内容多，每 10 页检测一次
                line_count = 0

                for page in doc:
                    if line_count >= 10000:
                        break
                    text = page.get_text()
                    if text:
                        for line in text.split("\n"):
                            if line_count >= 10000:
                                break
                            if line.strip():
                                chunk_parts.append(line.strip())
                                chunk_size += 1
                                line_count += 1

                    if chunk_size >= CHUNK_THRESHOLD:
                        chunk_text = "\n".join(chunk_parts)
                        should_stop, accumulated = matcher.scan_text_incremental(
                            chunk_text, accumulated
                        )
                        chunk_parts = []
                        chunk_size = 0
                        if should_stop:
                            is_sensitive, details = matcher.get_match_details(accumulated)
                            return is_sensitive, details, size_mb

                if chunk_parts:
                    chunk_text = "\n".join(chunk_parts)
                    should_stop, accumulated = matcher.scan_text_incremental(
                        chunk_text, accumulated
                    )

                is_sensitive, details = matcher.get_match_details(accumulated or {})
                return is_sensitive, details, size_mb
            finally:
                doc.close()
        except Exception as e:
            self.logger.warning(f"PyMuPDF 增量解析失败，回退: {file_path} - {e}")
            return self._read_pdf_pdfplumber_incremental(file_path, matcher, size_mb)

    def _read_pdf_pypdf(self, file_path):
        """使用 pypdf 解析 PDF（第三回退方案）"""
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            content = []
            line_count = 0
            total_chars = 0
            for page in reader.pages:
                if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                    break
                text = page.extract_text()
                if text:
                    for line in text.split("\n"):
                        if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                            break
                        if line.strip():
                            content.append(line.strip())
                            line_count += 1
                            total_chars += len(line)
            return "\n".join(content)
        except Exception as e:
            self.logger.warning(f"pypdf 解析失败: {file_path} - {e}")
            return ""

    def _read_pdf_pypdf2(self, file_path):
        """使用 PyPDF2 解析 PDF（最终回退方案，旧版包名）"""
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            content = []
            line_count = 0
            total_chars = 0
            for page in reader.pages:
                if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                    break
                text = page.extract_text()
                if text:
                    for line in text.split("\n"):
                        if line_count >= 10000 or total_chars >= MAX_TEXT_CHARS:
                            break
                        if line.strip():
                            content.append(line.strip())
                            line_count += 1
                            total_chars += len(line)
            return "\n".join(content)
        except Exception as e:
            self.logger.warning(f"PyPDF2 解析失败: {file_path} - {e}")
            return ""

    def _read_pdf_pdfplumber_incremental(self, file_path, matcher, size_mb):
        """使用 pdfplumber 增量提取+匹配（回退方案）"""
        try:
            import pdfplumber
            accumulated = None
            chunk_parts = []
            chunk_size = 0
            CHUNK_THRESHOLD = 5  # pdfplumber 慢，每 5 页检测一次
            line_count = 0

            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    if line_count >= 10000:
                        break
                    text = page.extract_text()
                    if text:
                        for line in text.split("\n"):
                            if line_count >= 10000:
                                break
                            if line.strip():
                                chunk_parts.append(line.strip())
                                chunk_size += 1
                                line_count += 1

                    if chunk_size >= CHUNK_THRESHOLD:
                        chunk_text = "\n".join(chunk_parts)
                        should_stop, accumulated = matcher.scan_text_incremental(
                            chunk_text, accumulated
                        )
                        chunk_parts = []
                        chunk_size = 0
                        if should_stop:
                            is_sensitive, details = matcher.get_match_details(accumulated)
                            return is_sensitive, details, size_mb

            if chunk_parts:
                chunk_text = "\n".join(chunk_parts)
                should_stop, accumulated = matcher.scan_text_incremental(
                    chunk_text, accumulated
                )

            is_sensitive, details = matcher.get_match_details(accumulated or {})
            return is_sensitive, details, size_mb
        except Exception as e:
            self.logger.warning(f"PDF 增量解析失败: {file_path} - {e}")
            return False, [], size_mb

    def _read_pdf_pypdf_incremental(self, file_path, matcher, size_mb):
        """使用 pypdf 增量提取+匹配"""
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            accumulated = None
            chunk_parts = []
            chunk_size = 0
            CHUNK_THRESHOLD = 10
            line_count = 0

            for page in reader.pages:
                if line_count >= 10000:
                    break
                text = page.extract_text()
                if text:
                    for line in text.split("\n"):
                        if line_count >= 10000:
                            break
                        if line.strip():
                            chunk_parts.append(line.strip())
                            chunk_size += 1
                            line_count += 1

                if chunk_size >= CHUNK_THRESHOLD:
                    chunk_text = "\n".join(chunk_parts)
                    should_stop, accumulated = matcher.scan_text_incremental(
                        chunk_text, accumulated
                    )
                    chunk_parts = []
                    chunk_size = 0
                    if should_stop:
                        is_sensitive, details = matcher.get_match_details(accumulated)
                        return is_sensitive, details, size_mb

            if chunk_parts:
                chunk_text = "\n".join(chunk_parts)
                should_stop, accumulated = matcher.scan_text_incremental(
                    chunk_text, accumulated
                )

            is_sensitive, details = matcher.get_match_details(accumulated or {})
            return is_sensitive, details, size_mb
        except Exception as e:
            self.logger.warning(f"pypdf 增量解析失败: {file_path} - {e}")
            return False, [], size_mb

    def _read_pdf_pypdf2_incremental(self, file_path, matcher, size_mb):
        """使用 PyPDF2 增量提取+匹配（旧版包名）"""
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            accumulated = None
            chunk_parts = []
            chunk_size = 0
            CHUNK_THRESHOLD = 10
            line_count = 0

            for page in reader.pages:
                if line_count >= 10000:
                    break
                text = page.extract_text()
                if text:
                    for line in text.split("\n"):
                        if line_count >= 10000:
                            break
                        if line.strip():
                            chunk_parts.append(line.strip())
                            chunk_size += 1
                            line_count += 1

                if chunk_size >= CHUNK_THRESHOLD:
                    chunk_text = "\n".join(chunk_parts)
                    should_stop, accumulated = matcher.scan_text_incremental(
                        chunk_text, accumulated
                    )
                    chunk_parts = []
                    chunk_size = 0
                    if should_stop:
                        is_sensitive, details = matcher.get_match_details(accumulated)
                        return is_sensitive, details, size_mb

            if chunk_parts:
                chunk_text = "\n".join(chunk_parts)
                should_stop, accumulated = matcher.scan_text_incremental(
                    chunk_text, accumulated
                )

            is_sensitive, details = matcher.get_match_details(accumulated or {})
            return is_sensitive, details, size_mb
        except Exception as e:
            self.logger.warning(f"PyPDF2 增量解析失败: {file_path} - {e}")
            return False, [], size_mb
